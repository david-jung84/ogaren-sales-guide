"""
오가렌 세일즈가이드 정적 사이트 빌더
- PPT → PNG 슬라이드 변환 (PowerPoint COM)
- 슬라이드 제목 추출 (python-pptx)
- index.html / session1.html / session2.html 생성

매월 갱신 워크플로우:
  1) part6_revised.js 수정 → node split_and_build.js
  2) python build_site.py
  3) git add . && git commit && git push
"""
from __future__ import annotations
import os
import re
import json
import shutil
import time
from pathlib import Path

import pythoncom
import win32com.client
from pptx import Presentation

ROOT = Path(__file__).parent
SITE = ROOT  # GitHub Pages는 repo root에서 서빙
ASSETS_S1 = SITE / "assets" / "session1"
ASSETS_S2 = SITE / "assets" / "session2"
ASSETS_DL = SITE / "assets" / "download"

S1_PPT = ROOT / "오가렌_세일즈가이드_세션1.pptx"
S2_PPT = ROOT / "오가렌_세일즈가이드_세션2.pptx"
FULL_PPT = ROOT / "오가렌_매니저_교육_세일즈_가이드_v2.1.pptx"

# 세션별 TOC 정의 (제목 + 슬라이드 범위)
# 슬라이드 번호는 1-indexed, 세션 PPT 내 기준
SESSION1_TOC = [
    {"id": "cover",  "label": "표지",                 "from": 1, "to": 1},
    {"id": "usage",  "label": "본 가이드 사용법",      "from": 2, "to": 2},
    {"id": "toc",    "label": "목차",                 "from": 3, "to": 3},
    {"id": "part1",  "label": "Part 1 · 오가렌과 3대 브랜드", "from": 4,  "to": 11},
    {"id": "part2",  "label": "Part 2 · 매트리스 마스터하기", "from": 12, "to": 19},
    {"id": "part3",  "label": "Part 3 · 침대 프레임 마스터하기", "from": 20, "to": 24},
    {"id": "part4",  "label": "Part 4 · 인증 & 안전성", "from": 25, "to": 37},
    {"id": "end",    "label": "마무리",               "from": 38, "to": 99},
]

SESSION2_TOC = [
    {"id": "cover",  "label": "표지",                 "from": 1, "to": 1},
    {"id": "usage",  "label": "본 가이드 사용법",      "from": 2, "to": 2},
    {"id": "toc",    "label": "목차",                 "from": 3, "to": 3},
    {"id": "part1",  "label": "Part 1 · 고객 응대 프로세스", "from": 4,  "to": 10},
    {"id": "part2",  "label": "Part 2 · 세일즈 스크립트 & 거절 처리", "from": 11, "to": 19},
    {"id": "part3",  "label": "Part 3 · TOP 매니저 케이스 & 모범 사례", "from": 20, "to": 28},
    {"id": "part4",  "label": "Part 4 · 운영 · KPI · 부록", "from": 29, "to": 33},
    {"id": "part5",  "label": "Part 5 · 배송 실무 가이드", "from": 34, "to": 37},
    {"id": "appendix", "label": "Appendix · 경쟁사 비교 자료", "from": 38, "to": 43},
    {"id": "end",    "label": "마무리",               "from": 44, "to": 99},
]


def ppt_export_pngs(ppt_path: Path, out_dir: Path, width: int = 1600) -> int:
    """PowerPoint COM을 통해 슬라이드를 PNG로 export. 반환: 슬라이드 개수."""
    out_dir.mkdir(parents=True, exist_ok=True)
    for f in out_dir.glob("*.png"):
        f.unlink()

    pythoncom.CoInitialize()
    try:
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.Visible = 1  # 일부 버전에서 헤드리스 모드가 거부됨
        try:
            deck = app.Presentations.Open(str(ppt_path), WithWindow=False)
        except Exception:
            deck = app.Presentations.Open(str(ppt_path))
        n = deck.Slides.Count
        height = int(width * 9 / 16)
        deck.SaveCopyAs(str(out_dir / "_export"), 18, EmbedTrueTypeFonts=False)  # 17=PNG, 18=GIF? — use Export
        # 위 SaveCopyAs는 ppt별 차이가 있어 Slide.Export 직접 호출이 더 안정적
        # 정리 후 Slide.Export 사용
        for tmp in out_dir.glob("_export*"):
            if tmp.is_dir():
                shutil.rmtree(tmp, ignore_errors=True)
            else:
                tmp.unlink(missing_ok=True)
        for i in range(1, n + 1):
            slide = deck.Slides(i)
            target = out_dir / f"slide_{i:02d}.png"
            slide.Export(str(target), "PNG", width, height)
        deck.Close()
        return n
    finally:
        try:
            app.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()


def extract_titles(ppt_path: Path) -> list[str]:
    """python-pptx로 각 슬라이드의 제목 텍스트 추출. 없으면 빈 문자열."""
    titles: list[str] = []
    pres = Presentation(str(ppt_path))
    for slide in pres.slides:
        title = ""
        # 1순위: title placeholder
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = (slide.shapes.title.text_frame.text or "").strip()
        if not title:
            # 2순위: 가장 위쪽의 큰 폰트 텍스트
            best = None
            for shp in slide.shapes:
                if not shp.has_text_frame:
                    continue
                txt = (shp.text_frame.text or "").strip()
                if not txt:
                    continue
                # 첫 줄만
                first_line = txt.splitlines()[0].strip()
                # 너무 긴 거(>40자) 제외
                if len(first_line) > 60:
                    continue
                key = (shp.top or 0, -(shp.width or 0))
                if best is None or key < best[0]:
                    best = (key, first_line)
            if best:
                title = best[1]
        # 정리
        title = re.sub(r"\s+", " ", title).strip()
        titles.append(title)
    return titles


# ============================================================
# HTML 생성
# ============================================================

CSS = """
@import url('https://fonts.googleapis.com/css2?family=Pretendard:wght@300;400;500;600;700;800;900&family=Cormorant+Garamond:ital,wght@0,400;0,500;0,600;0,700;1,400&display=swap');

* { box-sizing: border-box; margin: 0; padding: 0; }
:root {
  --bg: #0A0A0B;
  --bg-2: #111114;
  --surface: #15151A;
  --surface-2: #1C1C22;
  --border: rgba(255,255,255,0.08);
  --border-strong: rgba(255,255,255,0.14);
  --text: #F4F2EE;
  --text-2: #B8B5AE;
  --text-3: #76736D;
  --accent: #E8DCC4;
  --accent-warm: #C9B894;
  --serif: 'Cormorant Garamond', 'Noto Serif KR', 'Apple SD Gothic Neo', serif;
  --sans: 'Pretendard','Apple SD Gothic Neo','Malgun Gothic',system-ui,sans-serif;
}
html, body { font-family: var(--sans); color: var(--text); background: var(--bg); line-height: 1.6; -webkit-font-smoothing: antialiased; font-feature-settings: 'ss01','ss02'; }
body { background:
  radial-gradient(1200px 600px at 80% -10%, rgba(232,220,196,0.06), transparent 60%),
  radial-gradient(900px 500px at 0% 10%, rgba(201,184,148,0.04), transparent 60%),
  var(--bg);
  min-height: 100vh;
}
a { color: inherit; text-decoration: none; }
img { display: block; }
.container { max-width: 1320px; margin: 0 auto; padding: 0 32px; }
@media (max-width: 640px) { .container { padding: 0 20px; } }

/* ===== Index ===== */
.nav { padding: 28px 0; }
.nav .container { display: flex; justify-content: space-between; align-items: center; }
.nav .brand { font-family: var(--serif); font-size: 22px; font-weight: 500; letter-spacing: 0.5px; color: var(--accent); }
.nav .edition { font-size: 11px; letter-spacing: 2px; text-transform: uppercase; color: var(--text-3); }

.hero { padding: 80px 0 60px; }
.hero .container { position: relative; }
.hero .kicker { display: inline-block; font-size: 11px; letter-spacing: 3px; text-transform: uppercase; color: var(--accent-warm); margin-bottom: 28px; padding: 6px 14px; border: 1px solid var(--border-strong); border-radius: 999px; }
.hero h1 { font-family: var(--serif); font-weight: 500; font-size: clamp(46px, 8vw, 96px); line-height: 1.02; letter-spacing: -1.5px; margin-bottom: 28px; }
.hero h1 em { font-style: normal; color: var(--accent); font-weight: 400; }
.hero h1 .light { color: var(--text-2); font-weight: 300; }
.hero .lead { font-size: 17px; color: var(--text-2); max-width: 580px; line-height: 1.7; font-weight: 300; }
.hero .meta { margin-top: 36px; display: flex; gap: 36px; flex-wrap: wrap; font-size: 12px; letter-spacing: 1.5px; text-transform: uppercase; color: var(--text-3); }
.hero .meta b { color: var(--text); font-weight: 500; }

.divider { height: 1px; background: var(--border); margin: 24px 0; }

.sessions { padding: 60px 0 40px; }
.section-label { display: flex; justify-content: space-between; align-items: end; margin-bottom: 36px; gap: 24px; flex-wrap: wrap; }
.section-label h2 { font-family: var(--serif); font-size: 42px; font-weight: 500; letter-spacing: -0.5px; }
.section-label .num { font-size: 12px; letter-spacing: 2px; text-transform: uppercase; color: var(--text-3); }
.section-label .sub { font-size: 14px; color: var(--text-2); max-width: 360px; font-weight: 300; }

.session-grid { display: grid; grid-template-columns: 1fr 1fr; gap: 24px; }
@media (max-width: 860px) { .session-grid { grid-template-columns: 1fr; } }

.session-card {
  background: var(--surface);
  border: 1px solid var(--border);
  border-radius: 4px;
  padding: 48px 40px;
  position: relative;
  overflow: hidden;
  transition: border-color .35s ease, transform .35s ease, background .35s ease;
  display: block;
}
.session-card::before {
  content: ""; position: absolute; inset: 0;
  background: linear-gradient(135deg, rgba(232,220,196,0.06), transparent 50%);
  opacity: 0; transition: opacity .35s ease;
  pointer-events: none;
}
.session-card:hover { border-color: var(--accent-warm); transform: translateY(-2px); }
.session-card:hover::before { opacity: 1; }
.session-card .num { font-family: var(--serif); font-size: 14px; letter-spacing: 4px; color: var(--accent-warm); text-transform: uppercase; }
.session-card .num em { font-style: normal; }
.session-card h3 { font-family: var(--serif); font-size: 44px; font-weight: 500; line-height: 1.1; margin: 14px 0 18px; letter-spacing: -0.5px; }
.session-card h3 em { font-style: normal; color: var(--accent); font-weight: 400; }
.session-card .desc { font-size: 13px; color: var(--text-3); letter-spacing: 1.5px; text-transform: uppercase; margin-bottom: 28px; }
.session-card .parts { list-style: none; margin-top: 24px; border-top: 1px solid var(--border); padding-top: 24px; }
.session-card .parts li { font-size: 14px; color: var(--text-2); padding: 6px 0; display: flex; gap: 14px; align-items: baseline; font-weight: 300; }
.session-card .parts li b { font-family: var(--serif); font-style: normal; font-weight: 400; color: var(--accent-warm); font-size: 13px; width: 52px; flex-shrink: 0; letter-spacing: 1px; }
.session-card .footer { margin-top: 36px; display: flex; justify-content: space-between; align-items: center; padding-top: 24px; border-top: 1px solid var(--border); }
.session-card .count { font-size: 12px; letter-spacing: 1.5px; text-transform: uppercase; color: var(--text-3); }
.session-card .go { font-size: 13px; letter-spacing: 1.5px; text-transform: uppercase; color: var(--accent); display: inline-flex; align-items: center; gap: 10px; transition: gap .25s ease; }
.session-card:hover .go { gap: 16px; }

.downloads { padding: 60px 0 100px; }
.downloads h3 { font-family: var(--serif); font-size: 28px; font-weight: 500; margin-bottom: 24px; letter-spacing: -0.3px; }
.dl-grid { display: grid; grid-template-columns: repeat(3, 1fr); gap: 12px; }
@media (max-width: 720px) { .dl-grid { grid-template-columns: 1fr; } }
.dl-item { border: 1px solid var(--border); padding: 24px 24px; border-radius: 4px; transition: border-color .25s ease, background .25s ease; }
.dl-item:hover { border-color: var(--accent-warm); background: var(--surface); }
.dl-item .lbl { font-size: 11px; letter-spacing: 2px; text-transform: uppercase; color: var(--text-3); margin-bottom: 8px; }
.dl-item .name { font-size: 15px; color: var(--text); font-weight: 500; margin-bottom: 12px; }
.dl-item .arrow { color: var(--accent); font-size: 13px; letter-spacing: 1.5px; text-transform: uppercase; }

.footer { padding: 40px 0; border-top: 1px solid var(--border); margin-top: 60px; }
.footer .container { display: flex; justify-content: space-between; align-items: center; gap: 16px; flex-wrap: wrap; font-size: 12px; letter-spacing: 1px; color: var(--text-3); }
.footer .brand-mark { font-family: var(--serif); font-size: 14px; letter-spacing: 1px; color: var(--text-2); }

/* ===== Session viewer ===== */
.topbar { background: rgba(10,10,11,0.85); backdrop-filter: blur(20px); -webkit-backdrop-filter: blur(20px); border-bottom: 1px solid var(--border); padding: 18px 0; position: sticky; top: 0; z-index: 20; }
.topbar .container { display: flex; align-items: center; justify-content: space-between; gap: 16px; }
.topbar .back { font-size: 12px; letter-spacing: 1.5px; text-transform: uppercase; color: var(--text-2); display: inline-flex; align-items: center; gap: 10px; transition: color .25s ease; }
.topbar .back:hover { color: var(--accent); }
.topbar .title { font-family: var(--serif); font-size: 20px; font-weight: 500; letter-spacing: -0.3px; }
.topbar .title em { font-style: normal; color: var(--accent); }
.topbar .download { font-size: 11px; letter-spacing: 2px; text-transform: uppercase; padding: 10px 18px; border: 1px solid var(--border-strong); border-radius: 2px; transition: border-color .25s ease, color .25s ease, background .25s ease; color: var(--text); }
.topbar .download:hover { border-color: var(--accent-warm); color: var(--accent); background: var(--surface); }

.session-hero { padding: 80px 0 40px; }
.session-hero .kicker { font-size: 11px; letter-spacing: 3px; text-transform: uppercase; color: var(--accent-warm); margin-bottom: 24px; }
.session-hero h1 { font-family: var(--serif); font-size: clamp(40px, 6vw, 72px); font-weight: 500; line-height: 1.05; letter-spacing: -1px; }
.session-hero h1 em { font-style: normal; color: var(--accent); font-weight: 400; }

.layout { display: grid; grid-template-columns: 240px 1fr; gap: 56px; padding: 32px 0 80px; }
@media (max-width: 980px) { .layout { grid-template-columns: 1fr; gap: 24px; } }

.sidebar { position: sticky; top: 90px; align-self: start; max-height: calc(100vh - 110px); overflow-y: auto; padding-right: 12px; }
.sidebar::-webkit-scrollbar { width: 4px; }
.sidebar::-webkit-scrollbar-thumb { background: var(--border-strong); border-radius: 2px; }
.sidebar h4 { font-size: 10px; letter-spacing: 3px; text-transform: uppercase; color: var(--text-3); margin-bottom: 18px; }
.sidebar ul { list-style: none; }
.sidebar li a { display: block; padding: 12px 0; font-size: 13px; color: var(--text-2); border-bottom: 1px solid var(--border); transition: color .25s ease, padding-left .25s ease; line-height: 1.4; }
.sidebar li a:hover { color: var(--accent); padding-left: 8px; }
.sidebar li a .pnum { font-family: var(--serif); font-style: normal; color: var(--accent-warm); margin-right: 10px; font-size: 12px; letter-spacing: 1px; }

.main { min-width: 0; }
.part-section { margin-bottom: 80px; scroll-margin-top: 100px; }
.part-section .part-head { margin-bottom: 32px; display: flex; align-items: baseline; gap: 24px; flex-wrap: wrap; padding-bottom: 20px; border-bottom: 1px solid var(--border); }
.part-section .part-head .pnum { font-family: var(--serif); font-style: normal; font-size: 18px; color: var(--accent-warm); letter-spacing: 1px; }
.part-section .part-head h2 { font-family: var(--serif); font-size: 36px; font-weight: 500; letter-spacing: -0.5px; line-height: 1.1; flex: 1; min-width: 0; }
.part-section .part-head .range { font-size: 11px; letter-spacing: 2px; text-transform: uppercase; color: var(--text-3); }

.slide-grid { display: grid; grid-template-columns: repeat(2, 1fr); gap: 20px; }
@media (max-width: 700px) { .slide-grid { grid-template-columns: 1fr; } }
.slide-card { position: relative; background: var(--surface); border: 1px solid var(--border); border-radius: 4px; overflow: hidden; cursor: zoom-in; transition: border-color .35s ease, transform .35s ease; }
.slide-card:hover { border-color: var(--accent-warm); transform: translateY(-3px); }
.slide-card .thumb-wrap { position: relative; overflow: hidden; }
.slide-card .thumb { width: 100%; aspect-ratio: 16/9; object-fit: cover; background: #1a1a1f; transition: transform .6s cubic-bezier(.2,.6,.2,1); }
.slide-card:hover .thumb { transform: scale(1.04); }
.slide-card .num-badge { position: absolute; top: 14px; left: 14px; font-family: var(--serif); font-style: normal; font-size: 13px; color: var(--accent); background: rgba(10,10,11,0.7); backdrop-filter: blur(8px); padding: 4px 10px; border-radius: 2px; letter-spacing: 1px; }
.slide-card .meta { padding: 16px 18px; display: flex; justify-content: space-between; align-items: center; gap: 10px; border-top: 1px solid var(--border); }
.slide-card .meta b { color: var(--text); font-size: 13px; font-weight: 500; flex: 1; letter-spacing: 0.3px; }
.slide-card .meta .view { font-size: 10px; letter-spacing: 2px; text-transform: uppercase; color: var(--text-3); }

/* Lightbox */
.lb { position: fixed; inset: 0; background: rgba(8,8,10,0.96); backdrop-filter: blur(8px); display: none; align-items: center; justify-content: center; z-index: 100; padding: 32px; }
.lb.open { display: flex; }
.lb img { max-width: 100%; max-height: 88vh; box-shadow: 0 40px 100px rgba(0,0,0,0.7); border: 1px solid var(--border); }
.lb .close, .lb .nav { position: absolute; color: var(--text); background: rgba(255,255,255,0.06); border: 1px solid var(--border-strong); cursor: pointer; font-size: 18px; width: 48px; height: 48px; border-radius: 999px; display: flex; align-items: center; justify-content: center; transition: background .25s ease, border-color .25s ease; backdrop-filter: blur(10px); }
.lb .close:hover, .lb .nav:hover { background: rgba(232,220,196,0.1); border-color: var(--accent-warm); color: var(--accent); }
.lb .close { top: 28px; right: 28px; }
.lb .prev { left: 28px; top: 50%; transform: translateY(-50%); }
.lb .next { right: 28px; top: 50%; transform: translateY(-50%); }
.lb .counter { position: absolute; bottom: 32px; left: 50%; transform: translateX(-50%); font-family: var(--serif); font-style: normal; font-size: 14px; color: var(--text-2); letter-spacing: 2px; }
.lb .counter b { color: var(--accent); font-weight: 500; font-style: normal; }

/* Reveal on scroll */
.reveal { opacity: 0; transform: translateY(20px); transition: opacity .8s ease, transform .8s ease; }
.reveal.in { opacity: 1; transform: translateY(0); }
"""

LIGHTBOX_JS = """
(function(){
  const cards = Array.from(document.querySelectorAll('.slide-card'));
  if (cards.length) {
    const lb = document.createElement('div');
    lb.className = 'lb';
    lb.innerHTML = `
      <button class="close" aria-label="닫기">×</button>
      <button class="nav prev" aria-label="이전">‹</button>
      <img alt="" />
      <button class="nav next" aria-label="다음">›</button>
      <div class="counter"><b></b> <span class="sep">/</span> <span class="total"></span></div>
    `;
    document.body.appendChild(lb);
    const img = lb.querySelector('img');
    const cur = lb.querySelector('.counter b');
    const tot = lb.querySelector('.counter .total');
    let idx = 0;
    function show(i){
      idx = (i + cards.length) % cards.length;
      const im = cards[idx].querySelector('img');
      const src = im.getAttribute('data-full') || im.src;
      img.src = src;
      cur.textContent = String(idx+1).padStart(2,'0');
      tot.textContent = String(cards.length).padStart(2,'0');
    }
    function open(i){ show(i); lb.classList.add('open'); document.body.style.overflow = 'hidden'; }
    function close(){ lb.classList.remove('open'); document.body.style.overflow = ''; img.src=''; }
    cards.forEach((c, i) => c.addEventListener('click', () => open(i)));
    lb.querySelector('.close').onclick = close;
    lb.querySelector('.prev').onclick = (e) => { e.stopPropagation(); show(idx-1); };
    lb.querySelector('.next').onclick = (e) => { e.stopPropagation(); show(idx+1); };
    lb.addEventListener('click', (e) => { if (e.target === lb) close(); });
    document.addEventListener('keydown', (e) => {
      if (!lb.classList.contains('open')) return;
      if (e.key === 'Escape') close();
      if (e.key === 'ArrowLeft') show(idx-1);
      if (e.key === 'ArrowRight') show(idx+1);
    });
  }
  // reveal on scroll
  const obs = new IntersectionObserver((entries) => {
    entries.forEach(e => { if (e.isIntersecting) { e.target.classList.add('in'); obs.unobserve(e.target); } });
  }, { rootMargin: '-40px' });
  document.querySelectorAll('.reveal').forEach(el => obs.observe(el));
})();
"""


def html_index(s1_count: int, s2_count: int, updated: str) -> str:
    return f"""<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width, initial-scale=1" />
<title>오가렌 매장 매니저 교육 · 세일즈 가이드</title>
<meta name="description" content="오가렌 오프라인 쇼룸 매니저는 파는 사람이 아니라 수면 컨설턴트. 우리는 잘 잔다, 그래서 잘 안내한다." />
<style>{CSS}</style>
</head>
<body>
  <nav class="nav">
    <div class="container">
      <div class="brand">OGAREN</div>
      <div class="edition">Edition v2.1 · {updated}</div>
    </div>
  </nav>

  <section class="hero">
    <div class="container reveal">
      <span class="kicker">Manager Education Program</span>
      <h1>
        잘 자는 사람이<br/>
        <em>잠을 바꾼다.</em>
      </h1>
      <p class="lead">쇼룸 매니저는 파는 사람이 아니라 수면 컨설턴트. 한 사람의 잠을 바꾸는 매장 운영.</p>
      <div class="meta">
        <div><b>76</b> Slides</div>
        <div><b>2</b> Sessions</div>
        <div><b>9</b> Parts</div>
      </div>
    </div>
  </section>

  <section class="sessions">
    <div class="container">
      <div class="section-label reveal">
        <div>
          <div class="num">— The Two Sessions</div>
          <h2>두 개의 흐름.</h2>
        </div>
        <p class="sub">먼저 제품과 안전성으로 신뢰의 토대를 만들고, 그 다음 실제 응대·세일즈·운영으로 매장을 움직입니다.</p>
      </div>

      <div class="session-grid">
        <a class="session-card reveal" href="session1.html">
          <div class="num"><em>Session</em> 01</div>
          <h3>우리는<br/><em>잘 잔다.</em></h3>
          <div class="desc">Product · Certification · Safety</div>
          <ul class="parts">
            <li><b>Part 01</b>오가렌과 3대 브랜드</li>
            <li><b>Part 02</b>매트리스 마스터하기</li>
            <li><b>Part 03</b>침대 프레임 마스터하기</li>
            <li><b>Part 04</b>인증 & 안전성</li>
          </ul>
          <div class="footer">
            <span class="count">{s1_count} Slides</span>
            <span class="go">Enter →</span>
          </div>
        </a>

        <a class="session-card reveal" href="session2.html">
          <div class="num"><em>Session</em> 02</div>
          <h3>그래서<br/><em>잘 안내한다.</em></h3>
          <div class="desc">Consulting · Service · Operation</div>
          <ul class="parts">
            <li><b>Part 01</b>고객 응대 프로세스</li>
            <li><b>Part 02</b>세일즈 스크립트 & 거절 처리</li>
            <li><b>Part 03</b>TOP 매니저 케이스 (실명)</li>
            <li><b>Part 04</b>운영 · KPI · 부록</li>
            <li><b>Part 05</b>배송 실무 가이드</li>
          </ul>
          <div class="footer">
            <span class="count">{s2_count} Slides</span>
            <span class="go">Enter →</span>
          </div>
        </a>
      </div>
    </div>
  </section>

  <section class="downloads">
    <div class="container reveal">
      <div class="section-label">
        <div>
          <div class="num">— Originals</div>
          <h2>PPT로 가져가기.</h2>
        </div>
        <p class="sub">현장에서 화면에 띄우거나, 인쇄해서 두고 보기 위한 원본 파일.</p>
      </div>
      <div class="dl-grid">
        <a class="dl-item" href="assets/download/오가렌_세일즈가이드_세션1.pptx" download>
          <div class="lbl">Session 01 · PPTX</div>
          <div class="name">우리는 잘 잔다</div>
          <div class="arrow">↓ Download</div>
        </a>
        <a class="dl-item" href="assets/download/오가렌_세일즈가이드_세션2.pptx" download>
          <div class="lbl">Session 02 · PPTX</div>
          <div class="name">그래서 잘 안내한다</div>
          <div class="arrow">↓ Download</div>
        </a>
        <a class="dl-item" href="assets/download/오가렌_매니저_교육_세일즈_가이드_v2.1.pptx" download>
          <div class="lbl">Unified · v2.1 · 70 Slides</div>
          <div class="name">통합본 (전체)</div>
          <div class="arrow">↓ Download</div>
        </a>
      </div>
    </div>
  </section>

  <footer class="footer">
    <div class="container">
      <div class="brand-mark">OGAREN <em>Manager Education</em></div>
      <div>v2.1 · last updated {updated} · © DX팀</div>
    </div>
  </footer>
  <script>{LIGHTBOX_JS}</script>
</body>
</html>
"""


def html_session(session_name: str, session_label: str, subtitle: str, toc: list, titles: list[str], asset_dir: str, ppt_name: str) -> str:
    # 섹션별 슬라이드 정리
    parts_html: list[str] = []
    sidebar_html: list[str] = []
    # 사이드바 라벨용 파트 번호 매김 (Part 섹션만)
    part_idx = 0
    for sec in toc:
        slides_in_section = []
        for i in range(sec["from"], min(sec["to"], len(titles)) + 1):
            t = titles[i - 1] if i - 1 < len(titles) else ""
            slides_in_section.append((i, t))
        if not slides_in_section:
            continue
        cards = []
        for i, t in slides_in_section:
            t_display = t if t else f"슬라이드 {i}"
            cards.append(
                f'<div class="slide-card">'
                f'<div class="thumb-wrap">'
                f'<img class="thumb" loading="lazy" '
                f'src="{asset_dir}/slide_{i:02d}.png" '
                f'data-full="{asset_dir}/slide_{i:02d}.png" '
                f'alt="slide {i}" />'
                f'<span class="num-badge">{i:02d}</span>'
                f'</div>'
                f'<div class="meta"><b>{t_display}</b><span class="view">View</span></div>'
                f'</div>'
            )
        # 섹션 헤드 구성
        is_part = sec["id"].startswith("part")
        if is_part:
            part_idx += 1
            part_num = f"<span class='pnum'>Part {part_idx:02d}</span>"
            # label에서 "Part N · " 접두 제거
            heading = re.sub(r"^Part\s*\d+\s*[·\.]\s*", "", sec["label"])
        else:
            part_num = ""
            heading = sec["label"]
        parts_html.append(
            f'<section class="part-section reveal" id="{sec["id"]}">'
            f'<div class="part-head">{part_num}<h2>{heading}</h2>'
            f'<div class="range">{sec["from"]:02d} — {min(sec["to"], len(titles)):02d}</div></div>'
            f'<div class="slide-grid">{"".join(cards)}</div>'
            f'</section>'
        )
        # 사이드바
        if is_part:
            sidebar_html.append(
                f'<li><a href="#{sec["id"]}"><span class="pnum">{part_idx:02d}</span>{heading}</a></li>'
            )
        else:
            sidebar_html.append(
                f'<li><a href="#{sec["id"]}">{heading}</a></li>'
            )

    # 타이틀 파싱: "세션 1" → "Session 01", subtitle "우리는 잘 잔다" → 마지막 단어 italic
    session_num = session_label.replace("세션 ", "")
    subtitle_em = re.sub(r"(잘\s*잔다|잘\s*판다|잘\s*안내한다|잠을\s*바꾼다)\.?", r"<em>\1.</em>", subtitle)

    return f"""<!doctype html>
<html lang="ko">
<head>
<meta charset="utf-8" />
<meta name="viewport" content="width=device-width, initial-scale=1" />
<title>{session_label} · {subtitle} | 오가렌 매니저 교육</title>
<meta name="description" content="오가렌 매니저 교육 · {session_label} · {subtitle}" />
<style>{CSS}</style>
</head>
<body>
  <div class="topbar">
    <div class="container">
      <a href="index.html" class="back">← Back to Index</a>
      <div class="title"><em>Session</em> {session_num} · {subtitle}</div>
      <a href="assets/download/{ppt_name}" class="download" download>↓ PPTX</a>
    </div>
  </div>

  <section class="session-hero">
    <div class="container reveal">
      <div class="kicker">Session {session_num}</div>
      <h1>{subtitle_em}</h1>
    </div>
  </section>

  <div class="container">
    <div class="layout">
      <aside class="sidebar">
        <h4>— Contents</h4>
        <ul>{"".join(sidebar_html)}</ul>
      </aside>
      <main class="main">
        {"".join(parts_html)}
      </main>
    </div>
  </div>

  <footer class="footer">
    <div class="container">
      <div class="brand-mark">OGAREN <em>Manager Education</em></div>
      <div>v2.1 · {session_label} · © DX팀</div>
    </div>
  </footer>
  <script>{LIGHTBOX_JS}</script>
</body>
</html>
"""


def main():
    import sys
    skip_png = "--no-png" in sys.argv
    SITE.mkdir(exist_ok=True)
    ASSETS_S1.mkdir(parents=True, exist_ok=True)
    ASSETS_S2.mkdir(parents=True, exist_ok=True)
    ASSETS_DL.mkdir(parents=True, exist_ok=True)

    if skip_png:
        s1_count = len(list(ASSETS_S1.glob("slide_*.png")))
        s2_count = len(list(ASSETS_S2.glob("slide_*.png")))
        print(f"[1-2/4] PNG export skipped (--no-png) - reusing {s1_count}+{s2_count} slides")
    else:
        print("[1/4] session1 PPT -> PNG ...")
        s1_count = ppt_export_pngs(S1_PPT, ASSETS_S1)
        print(f"    {s1_count} slides")

        print("[2/4] session2 PPT -> PNG ...")
        s2_count = ppt_export_pngs(S2_PPT, ASSETS_S2)
        print(f"    {s2_count} slides")

    print("[3/4] 슬라이드 제목 추출 …")
    s1_titles = extract_titles(S1_PPT)
    s2_titles = extract_titles(S2_PPT)

    print("[4/4] HTML 빌드 + PPT 사본 복사 …")
    for src in [S1_PPT, S2_PPT, FULL_PPT]:
        if src.exists():
            shutil.copy2(src, ASSETS_DL / src.name)

    updated = time.strftime("%Y-%m-%d")

    (SITE / "index.html").write_text(
        html_index(s1_count, s2_count, updated), encoding="utf-8"
    )
    (SITE / "session1.html").write_text(
        html_session("session1", "세션 1", "우리는 잘 잔다",
                     SESSION1_TOC, s1_titles, "assets/session1",
                     S1_PPT.name),
        encoding="utf-8",
    )
    (SITE / "session2.html").write_text(
        html_session("session2", "세션 2", "그래서 잘 안내한다",
                     SESSION2_TOC, s2_titles, "assets/session2",
                     S2_PPT.name),
        encoding="utf-8",
    )

    # 메타 JSON (디버그용)
    (SITE / "_meta.json").write_text(
        json.dumps({
            "updated": updated,
            "session1": {"count": s1_count, "titles": s1_titles},
            "session2": {"count": s2_count, "titles": s2_titles},
        }, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )

    print(f"\nDONE -> ./index.html")
    print(f"  session1: {s1_count} slides")
    print(f"  session2: {s2_count} slides")


if __name__ == "__main__":
    main()
