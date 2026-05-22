# -*- coding: utf-8 -*-
"""
세션 3 PDF/PPTX → PNG 슬라이드 렌더링.
세션 1/2와 동일한 구조: assets/session3/slide_NNN.png

산출:
  - assets/session3/slide_001.png ... slide_NNN.png
  - _session3_slides.json (TOC + 슬라이드 메타)
"""
from __future__ import annotations
import io
import json
import os
import shutil
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image

ROOT = Path(__file__).parent
DL = ROOT / "assets" / "download" / "session3"
OUT = ROOT / "assets" / "session3"
OUT.mkdir(parents=True, exist_ok=True)
# 기존 slide_*.png 청소
for f in OUT.glob("slide_*.png"):
    f.unlink()

W, H = 1600, 900  # 16:9 PNG 크기

# Part 정의: (id, label, source_pdf, slide_titles)
# slide_titles는 길이가 source PDF 페이지 수와 같아야 함.
PARTS = [
    # ─ 0 표지 (자동 생성된 1슬라이드, 세션 표제)
    # 매트리스 라인업 4종
    {
        "id": "signature",
        "label": "Part 01 · 시그니처 매트리스",
        "src": "01_signature_lineup_251216.pdf",
        "titles": [
            "SIGNATURE · 표지", "라인업 소개", "포지셔닝", "시그니처 1 (포근)",
            "시그니처 1 · 디테일", "시그니처 1 · 구성", "시그니처 2 (균형)",
            "시그니처 2 · 디테일", "시그니처 2 · 구성", "시그니처 3 (탄탄)",
            "시그니처 3 · 디테일", "시그니처 3 · 구성", "3종 비교 1",
            "3종 비교 2", "프레임 호환", "사이즈·가격", "케어", "FAQ",
            "체험 가이드", "CLOSING",
        ],
    },
    {
        "id": "hybrid",
        "label": "Part 02 · 하이브리드 매트리스",
        "src": "02_hybrid_lineup_251216.pdf",
        "titles": [
            "HYBRID · 표지", "라인업 소개", "포지셔닝", "구조 (포켓+메모리폼)",
            "소재", "라인업", "사이즈·가격", "케어", "FAQ", "체험 가이드", "CLOSING",
        ],
    },
    {
        "id": "stable",
        "label": "Part 03 · 스테이블 매트리스",
        "src": "03_stable_lineup_251216.pdf",
        "titles": [
            "STABLE · 표지", "라인업 소개", "포지셔닝", "구조 · 지지력",
            "소재", "라인업", "사이즈·가격", "케어", "FAQ", "체험 가이드", "CLOSING",
        ],
    },
    {
        "id": "balanced",
        "label": "Part 04 · 밸런스드 매트리스",
        "src": "04_balanced_lineup_251216.pdf",
        "titles": [
            "BALANCED · 표지", "라인업 소개", "포지셔닝", "구조 · 균형감",
            "소재", "라인업", "사이즈·가격", "케어", "FAQ", "체험 가이드", "CLOSING",
        ],
    },
    # 경쟁사 비교
    {
        "id": "compete",
        "label": "Part 05 · 경쟁사 비교 (3.3.3)",
        "src": "05_competitor_compare_251212.pdf",
        "titles": [
            "경쟁사 비교 · 표지",
            "12종 시리즈 특장 + 경도",
            "시그니처 vs 시몬스 · 씰리",
            "하이브리드 vs 코웨이 · 씰리 · 베스트슬립",
            "스테이블 vs 씰리 · 코웨이 · 에이스 · 한셀 · 일룸",
            "밸런스드 vs 에이스 · 코웨이 · 일룸",
        ],
    },
    # 브랜드 · 분류
    {
        "id": "nooer",
        "label": "Part 06 · 누어 교육자료",
        "src": "06_nooer_education_240923_lite.pdf",
        "titles": ["NOOER · 포근히, 탄탄히 — 브랜드 1장"],
    },
    {
        "id": "snc",
        "label": "Part 07 · 슬립퍼 vs 누어",
        "src": "07_sleeper_nooer_compare_240220_lite.pdf",
        "titles": ["SLEEPER & NOOER 비교 — 브랜드 포지셔닝"],
    },
    {
        "id": "classify",
        "label": "Part 08 · 슬립퍼 매트리스 분류표",
        "src": "09_mattress_classification_240513_lite.pdf",
        "titles": [
            "분류표 · 표지",
            "포지션 매트릭스 (포근 ~ 탄탄)",
            "라인업 한눈에",
            "소재 분류",
            "사이즈·가격 분류",
            "가격대별",
            "추천 가이드",
        ],
    },
    # 안전성
    {
        "id": "radon",
        "label": "Part 09 · 라돈 안전성",
        "src": "08_radon_safety_240220_lite.pdf",
        "titles": ["라돈 검사 데이터 + 인증 1장"],
    },
    # 베개
    {
        "id": "pillow",
        "label": "Part 10 · 경추베개",
        "src": "10_cervical_pillow_240214.pdf",
        "titles": ["경추베개 라인업 — 매장 응대용 1장"],
    },
]


def render_page_to_png(page: fitz.Page, out_path: Path):
    """16:9 PNG로 렌더링. 종횡비 다르면 검은 패딩."""
    r = page.rect
    aspect = r.width / r.height
    want = W / H
    # zoom to fit
    if aspect > want:
        zoom = W / r.width
    else:
        zoom = H / r.height
    pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
    img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
    # 16:9 캔버스 위에 가운데 정렬
    canvas = Image.new("RGB", (W, H), (15, 15, 18))
    x = (W - img.width) // 2
    y = (H - img.height) // 2
    canvas.paste(img, (x, y))
    canvas.save(out_path, "PNG", optimize=True)


# 슬라이드 1: 세션 표지 (자동 생성)
COVER_TXT_LARGE = "Session 03"
COVER_TXT_TITLE = "쇼룸 · 백화점 교육자료"
COVER_SUB = "Product Library · Lineup · Reference"

try:
    from PIL import ImageDraw, ImageFont
    cover = Image.new("RGB", (W, H), (10, 10, 11))
    draw = ImageDraw.Draw(cover)
    # 폰트 fallback
    def _font(size, bold=False):
        for name in ([
            "C:/Windows/Fonts/malgunbd.ttf" if bold else "C:/Windows/Fonts/malgun.ttf",
        ]):
            try:
                return ImageFont.truetype(name, size)
            except Exception:
                pass
        return ImageFont.load_default()
    f_small = _font(22)
    f_title = _font(96, bold=True)
    f_sub = _font(20)
    draw.text((80, 240), COVER_TXT_LARGE.upper(), fill=(201, 184, 148), font=f_small)
    draw.text((80, 290), COVER_TXT_TITLE, fill=(232, 220, 196), font=f_title)
    draw.text((80, 480), COVER_SUB, fill=(118, 115, 109), font=f_sub)
    # 가장자리 선
    draw.line([(80, 590), (W - 80, 590)], fill=(40, 40, 46), width=1)
    draw.text((80, 614), "OGAREN · Manager Education", fill=(118, 115, 109), font=f_sub)
    cover.save(OUT / "slide_001.png", "PNG", optimize=True)
    cover_ok = True
except Exception as e:
    print(f"cover 생성 실패: {e}")
    cover_ok = False

slide_idx = 2 if cover_ok else 1
toc = []
if cover_ok:
    toc.append({
        "id": "cover", "label": "표지",
        "from": 1, "to": 1,
        "slides": [{"num": 1, "title": "SESSION 03 · COVER"}],
    })

for part in PARTS:
    src = DL / part["src"]
    if not src.exists():
        print(f"[SKIP] {src.name} 없음")
        continue
    doc = fitz.open(src)
    n = doc.page_count
    titles = part["titles"]
    if len(titles) < n:
        titles = titles + [f"{titles[0] if titles else part['label']} · {i+1}" for i in range(len(titles), n)]
    elif len(titles) > n:
        titles = titles[:n]

    start = slide_idx
    slide_list = []
    for i in range(n):
        page = doc.load_page(i)
        out_p = OUT / f"slide_{slide_idx:03d}.png"
        render_page_to_png(page, out_p)
        slide_list.append({"num": slide_idx, "title": titles[i]})
        slide_idx += 1
    doc.close()
    toc.append({
        "id": part["id"],
        "label": part["label"],
        "from": start,
        "to": slide_idx - 1,
        "slides": slide_list,
    })
    print(f"  {part['id']}: slides {start} - {slide_idx - 1}  ({n}p)")

# 하디 PPTX는 PowerPoint COM으로 따로 렌더
pptx_path = DL / "11_hardy_bed_education.pptx"
if pptx_path.exists():
    try:
        import pythoncom
        import win32com.client
        from pptx import Presentation

        pythoncom.CoInitialize()
        ppt_app = win32com.client.Dispatch("PowerPoint.Application")
        ppt_app.Visible = 1
        ppt_app.DisplayAlerts = 1

        tmp_dir = OUT / "_hardy_tmp"
        if tmp_dir.exists():
            shutil.rmtree(tmp_dir)
        tmp_dir.mkdir()

        deck = ppt_app.Presentations.Open(str(pptx_path), WithWindow=False)
        # 슬라이드 제목 추출
        prs = Presentation(str(pptx_path))
        hardy_titles = []
        for i, sl in enumerate(prs.slides):
            t = None
            for shp in sl.shapes:
                if shp.has_text_frame and shp.text_frame.text.strip():
                    t = shp.text_frame.text.strip().split("\n")[0][:40]
                    break
            hardy_titles.append(t or f"하디 · {i+1}")
        # 통째로 PNG export
        deck.SaveAs(str(tmp_dir.resolve()), 18)  # 18 = ppSaveAsPNG
        deck.Close()
        ppt_app.Quit()
        pythoncom.CoUninitialize()

        # SaveAs는 폴더에 Slide1.PNG, Slide2.PNG... 생성
        # 또는 같은 이름 폴더에 들어갈 수 있음 — 두 경우 모두 처리
        candidate_dirs = [tmp_dir] + [p for p in tmp_dir.iterdir() if p.is_dir()]
        png_files = []
        for d in candidate_dirs:
            png_files.extend(sorted(d.glob("Slide*.PNG")) + sorted(d.glob("슬라이드*.PNG")))
        png_files = sorted(set(png_files), key=lambda p: int("".join(c for c in p.stem if c.isdigit()) or "0"))

        start = slide_idx
        slide_list = []
        for i, src_png in enumerate(png_files):
            out_p = OUT / f"slide_{slide_idx:03d}.png"
            # 16:9 캔버스로 정규화
            img = Image.open(src_png).convert("RGB")
            ratio = img.width / img.height
            want = W / H
            if ratio > want:
                new_w = W; new_h = int(W / ratio)
            else:
                new_h = H; new_w = int(H * ratio)
            img = img.resize((new_w, new_h), Image.LANCZOS)
            canvas = Image.new("RGB", (W, H), (15, 15, 18))
            canvas.paste(img, ((W - new_w) // 2, (H - new_h) // 2))
            canvas.save(out_p, "PNG", optimize=True)
            title = hardy_titles[i] if i < len(hardy_titles) else f"하디 · {i+1}"
            slide_list.append({"num": slide_idx, "title": title})
            slide_idx += 1
        shutil.rmtree(tmp_dir)
        if slide_list:
            toc.append({
                "id": "hardy",
                "label": "Part 11 · 하디 침대",
                "from": start,
                "to": slide_idx - 1,
                "slides": slide_list,
            })
            print(f"  hardy: slides {start} - {slide_idx - 1}  ({len(slide_list)}p)")
    except Exception as e:
        print(f"하디 PPTX 렌더 실패: {e}")

total = slide_idx - 1
print(f"\nTOTAL: {total} slides")

(ROOT / "_session3_slides.json").write_text(
    json.dumps({"total": total, "toc": toc}, ensure_ascii=False, indent=2),
    encoding="utf-8",
)
print(f"OK toc -> _session3_slides.json")
